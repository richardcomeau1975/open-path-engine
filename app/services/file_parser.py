"""
File parser service.
Extracts text from uploaded student materials.
Supports: PDF, PPTX, DOCX, XLSX, TXT, MD
"""

import io
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document
from openpyxl import load_workbook


def parse_pdf(file_bytes: bytes) -> str:
    """Extract text from a PDF file."""
    reader = PdfReader(io.BytesIO(file_bytes))
    pages = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            pages.append(text.strip())
    return "\n\n".join(pages)


def parse_pptx(file_bytes: bytes) -> str:
    """Extract text from a PowerPoint file. Preserves slide order."""
    prs = Presentation(io.BytesIO(file_bytes))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = paragraph.text.strip()
                    if line:
                        texts.append(line)
        if texts:
            slides.append(f"--- Slide {i} ---\n" + "\n".join(texts))
    return "\n\n".join(slides)


def parse_docx(file_bytes: bytes) -> str:
    """Extract text from a Word document."""
    doc = Document(io.BytesIO(file_bytes))
    paragraphs = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            paragraphs.append(text)
    return "\n\n".join(paragraphs)


def parse_xlsx(file_bytes: bytes) -> str:
    """Extract text from an Excel file. Each sheet as a section."""
    wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    sheets = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            line = " | ".join(cells).strip(" |")
            if line:
                rows.append(line)
        if rows:
            sheets.append(f"--- Sheet: {sheet_name} ---\n" + "\n".join(rows))
    wb.close()
    return "\n\n".join(sheets)


def parse_text(file_bytes: bytes) -> str:
    """Extract text from a plain text or markdown file."""
    return file_bytes.decode("utf-8", errors="replace").strip()


# Map of file extension (lowercase, with dot) to parser function
PARSERS = {
    ".pdf": parse_pdf,
    ".pptx": parse_pptx,
    ".docx": parse_docx,
    ".xlsx": parse_xlsx,
    ".txt": parse_text,
    ".md": parse_text,
}


def parse_file(filename: str, file_bytes: bytes) -> str:
    """
    Parse a file and return extracted text.
    Raises ValueError if file type is not supported.
    """
    ext = "." + filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    parser = PARSERS.get(ext)
    if not parser:
        raise ValueError(f"Unsupported file type: {ext}")
    return parser(file_bytes)


def parse_multiple_files(files: list[tuple[str, bytes]]) -> str:
    """
    Parse multiple files and return concatenated text.
    files: list of (filename, file_bytes) tuples.
    Skips unsupported files with a warning header.
    """
    sections = []
    for filename, file_bytes in files:
        try:
            text = parse_file(filename, file_bytes)
            sections.append(f"=== {filename} ===\n{text}")
        except ValueError:
            sections.append(f"=== {filename} === [SKIPPED: unsupported file type]")
    return "\n\n".join(sections)
